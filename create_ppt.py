"""Generate Campus Cafe presentation as .pptx"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
BLUE_DARK  = RGBColor(0x1a, 0x36, 0x5d)
BLUE       = RGBColor(0x2b, 0x6c, 0xb0)
BLUE_MID   = RGBColor(0x31, 0x82, 0xce)
BLUE_LIGHT = RGBColor(0xeb, 0xf4, 0xff)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GRAY_700   = RGBColor(0x4a, 0x56, 0x68)
GRAY_500   = RGBColor(0x71, 0x80, 0x96)
BLACK      = RGBColor(0x1a, 0x20, 0x2c)
CODE_BG    = RGBColor(0xf7, 0xfa, 0xfc)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_font(run, size=18, bold=False, color=BLACK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def add_textbox(slide, x, y, w, h, text="", size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, name="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p.runs[0], size, bold, color, name)
    return tb


def add_bullet_list(slide, x, y, w, h, items, size=16, color=GRAY_700):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.space_after = Pt(6)
        p.level = 0
        for run in p.runs:
            set_font(run, size, color=color)
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = pPr.find(qn('a:buChar'))
        if buChar is None:
            buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')
        # indent
        pPr.set('marL', str(Emu(Inches(0.3))))
        pPr.set('indent', str(Emu(Inches(-0.2))))
    return tb


def add_code_box(slide, x, y, w, h, code_text, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = code_text
    set_font(p.runs[0], font_size, color=BLACK, name="Courier New")
    p.line_spacing = Pt(font_size * 1.6)
    return shape


def add_info_card(slide, x, y, w, h, title, body_items):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_LIGHT
    shape.line.fill.background()
    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.1), Inches(0.06), h - Inches(0.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), w - Inches(0.3), Inches(0.35),
                title, size=14, bold=True, color=BLUE_DARK)
    # Bullets
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.45), w - Inches(0.3), h - Inches(0.5),
                    body_items, size=12, color=GRAY_700)


def slide_header(slide, title, subtitle=""):
    add_bg_rect(slide, 0, 0, W, Inches(1.2), BLUE_DARK)
    add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.5),
                title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.6), Inches(0.7), Inches(10), Inches(0.4),
                    subtitle, size=14, color=RGBColor(0xbb, 0xcc, 0xee))
    # Footer
    add_textbox(slide, Inches(0.6), H - Inches(0.4), Inches(3), Inches(0.3),
                "Campus Cafe", size=10, color=GRAY_500)


# ======================== SLIDE 1: Title ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
# Gradient overlay - second rect
add_bg_rect(slide, Inches(6), 0, Inches(7.333), H, BLUE)
add_textbox(slide, 0, Inches(2.0), W, Inches(1.0),
            "CAMPUS CAFE", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.1), W, Inches(0.5),
            "A Full-Stack Point-of-Sale System for Campus Dining", size=20, color=RGBColor(0xcc, 0xdd, 0xee), align=PP_ALIGN.CENTER)
# Divider line
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.03))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = RGBColor(0x66, 0x88, 0xbb)
line_shape.line.fill.background()
add_textbox(slide, 0, Inches(4.3), W, Inches(0.8),
            "NEU 5001 — Final Project\nBuilt with Python  •  FastAPI  •  React", size=15, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ======================== SLIDE 2: Agenda ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Agenda")
items = [
    "Project Overview — What we built and why",
    "System Architecture — How the pieces fit together",
    "The Main Function — CLI application flow and thought process",
    "Class Design Deep Dive — Menu, Order, Inventory, Customer, Staff",
    "Data Persistence — Excel-based tracking system (Special Feature)",
    "Frontend Application — React web interface",
    "Live Demo & Q&A",
]
add_bullet_list(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(5), items, size=18, color=GRAY_700)


# ======================== SLIDE 3: Project Overview ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Project Overview", "What is Campus Cafe?")
add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
            "Campus Cafe is a point-of-sale (POS) system for a campus dining environment, supporting customer ordering and staff inventory management. It runs on the terminal AND as a modern web application.",
            size=15, color=GRAY_700)
# 4 cards
cw = Inches(5.8)
ch = Inches(2.0)
add_info_card(slide, Inches(0.6), Inches(2.4), cw, ch, "Customer Features",
              ["Browse menu by category (Salads, Drinks, Desserts, Sandwiches)",
               "Add / remove items from cart with quantity control",
               "Daily special with automatic 20% discount",
               "Checkout with loyalty points & printed receipt"])
add_info_card(slide, Inches(6.8), Inches(2.4), cw, ch, "Staff Features",
              ["Secure login with ID + password",
               "View real-time inventory levels",
               "Restock items in bulk",
               "Session summary on quit (all orders + total revenue)"])
add_info_card(slide, Inches(0.6), Inches(4.7), cw, ch, "Dual Interface",
              ["Terminal-based CLI application (main.py)",
               "React web application (frontend/)",
               "Both interfaces share the exact same backend classes"])
add_info_card(slide, Inches(6.8), Inches(4.7), cw, ch, "Data Persistence",
              ["All data stored in Excel (.xlsx) files",
               "Orders, customers, inventory, staff records",
               "Loyalty points accumulate across sessions"])


# ======================== SLIDE 4: Architecture ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "System Architecture", "How the pieces fit together")

# Architecture boxes
boxes = [
    ("Terminal CLI", "main.py", Inches(0.3)),
    ("Backend Classes", "Menu • Order • Inventory\nCustomer • Staff", Inches(3.3)),
    ("FastAPI Server", "app.py (REST API)", Inches(7.0)),
    ("React Frontend", "Vite + React Router", Inches(10.2)),
]
for title, desc, x_pos in boxes:
    is_center = "Backend" in title
    bg = BLUE if is_center else BLUE_LIGHT
    tc = WHITE if is_center else BLUE_DARK
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(2.0), Inches(2.5), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = BLUE_MID
    shape.line.width = Pt(2)
    add_textbox(slide, x_pos + Inches(0.15), Inches(2.1), Inches(2.2), Inches(0.35),
                title, size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
    add_textbox(slide, x_pos + Inches(0.15), Inches(2.45), Inches(2.2), Inches(0.6),
                desc, size=10, color=tc if is_center else GRAY_500, align=PP_ALIGN.CENTER)

# Arrows between boxes
for ax in [Inches(2.8), Inches(5.8), Inches(9.5)]:
    add_textbox(slide, ax, Inches(2.3), Inches(0.5), Inches(0.5),
                "→", size=24, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

# Excel box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.3), Inches(4.0), Inches(6.2), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE_LIGHT
shape.line.color.rgb = BLUE_MID
shape.line.width = Pt(2)
add_textbox(slide, Inches(3.5), Inches(4.05), Inches(5.8), Inches(0.3),
            "Excel Data Layer", size=13, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.5), Inches(4.35), Inches(5.8), Inches(0.4),
            "Menu.xlsx  •  Orders.xlsx  •  Inventory.xlsx  •  Customers.xlsx  •  Staffs.xlsx",
            size=10, color=GRAY_500, align=PP_ALIGN.CENTER)
# Arrow down
add_textbox(slide, Inches(6.2), Inches(3.25), Inches(0.5), Inches(0.5),
            "↓", size=24, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

# Design principle box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.4), Inches(12), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(5.5), Inches(11.4), Inches(0.3),
            "Design Principle", size=14, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.85), Inches(11.4), Inches(0.5),
            "The backend classes are shared between the CLI and the web app. The FastAPI server wraps the same Python classes that main.py uses, ensuring consistent behavior across both interfaces.",
            size=13, color=GRAY_700)


# ======================== SLIDE 5: Section - Main Function ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
add_textbox(slide, 0, Inches(1.8), W, Inches(0.4),
            "SECTION 1", size=16, bold=True, color=RGBColor(0x66, 0x88, 0xbb), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(2.5), W, Inches(0.8),
            "The Main Function", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.5), W, Inches(0.4),
            "CLI application flow and thought process", size=18, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ======================== SLIDE 6: Main Function - Code ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Main Function — Application Flow", "main.py: the entry point that orchestrates everything")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
            "Thought Process: We designed main.py as a menu-driven loop. The user sees numbered options, picks one, and the program delegates to a dedicated function. This keeps each feature isolated and the main loop clean.",
            size=14, color=GRAY_700)

# Left code block
add_textbox(slide, Inches(0.6), Inches(2.1), Inches(3), Inches(0.3),
            "ENTRY POINT & GLOBAL STATE", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(0.6), Inches(2.4), Inches(5.8), Inches(2.6),
"""from backend.menu import MenuItem, Menu
from backend.order import Order
from backend.inventory_ShuranRyu import Inventory
from backend.customer import Customer, Customers
from backend.staff import Staff, Staffs

MENU  = Menu()        # Load menu from Excel
order = Order(MENU)   # Current order
staffs = Staffs()     # Load staff credentials
session_orders = {}   # Track all orders""", font_size=11)

# Right code block
add_textbox(slide, Inches(6.8), Inches(2.1), Inches(3), Inches(0.3),
            "MAIN_MENU() — THE DRIVING LOOP", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.8), Inches(2.4), Inches(5.8), Inches(2.6),
"""def main_menu():
    while True:
        # Display options to user
        choice = input("Enter your choice: ")

        if   choice == "1": view_menu()
        elif choice == "2": order_item()
        elif choice == "3": delete_item()
        elif choice == "4": view_cart()
        elif choice == "5": checkout()
        elif choice == "s": login_staff_system()""", font_size=11)

# Bottom: key decision box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.3), Inches(12), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(5.4), Inches(11.4), Inches(0.3),
            "Why This Design?", size=14, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.75), Inches(11.4), Inches(0.7),
            "Each function (view_menu, order_item, checkout, etc.) is self-contained with its own input loop and validation. Global state (MENU, order, session_orders) is initialized once at the top and shared across functions. This mirrors how a real POS works — one session, one shared order, one menu.",
            size=13, color=GRAY_700)


# ======================== SLIDE 7: Ordering & Checkout ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Main Function — Ordering & Checkout", "Input validation, stock checking, and receipt generation")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(3), Inches(0.3),
            "ORDER_ITEM() — WITH STOCK VALIDATION", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.8),
"""def order_item():
    item = MENU.get_item(item_name)
    # Validate quantity (try/except for int)

    inventory = Inventory()
    counter = 0
    for i in range(num_int):
        if inventory.check_stock(item):
            counter += 1
            order.add_to_order(item)
            inventory.deduct_stock(item)
            inventory.save_to_excel()
        else:
            print("Out of stock")
            break
    print(f"You have ordered {counter} {item.name}.")""", font_size=11)

add_textbox(slide, Inches(6.8), Inches(1.4), Inches(3), Inches(0.3),
            "CHECKOUT() — RECEIPT & LOYALTY POINTS", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.8),
"""def checkout():
    subtotal = order.get_total()
    total_before_tax = order.get_total_after_discount()
    discount = subtotal - total_before_tax
    tax = order.TAX_RATE * total_before_tax
    total = total_before_tax + tax

    # Loyalty points
    customers = Customers()
    if customers.check_customer(telephone):
        customer = customers.get_customer(telephone)
        customer.earn_points(total)
    else:
        customer = Customer(telephone)
        customer.earn_points(total)
        customers.add_customer(customer)

    # Save & track session
    order.save_to_excel()
    session_orders[len(session_orders)] = {...}""", font_size=11)

# Key decision boxes
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.8), Inches(5.8), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(4.9), Inches(5.2), Inches(0.3),
            "Key Decision: Per-Unit Stock Check", size=13, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.25), Inches(5.2), Inches(1.0),
            "We check stock per unit in a loop rather than checking total availability upfront. This enables partial fulfillment — if 3 are requested but only 2 remain, the customer gets 2 instead of 0.",
            size=12, color=GRAY_700)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(4.8), Inches(5.8), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(7.1), Inches(4.9), Inches(5.2), Inches(0.3),
            "Key Decision: Auto Customer Registration", size=13, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(7.1), Inches(5.25), Inches(5.2), Inches(1.0),
            "At checkout, we check if the phone number exists. Returning customers earn more points; new customers are automatically registered. No signup flow needed — frictionless loyalty.",
            size=12, color=GRAY_700)


# ======================== SLIDE 8: Section - Class Design ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
add_textbox(slide, 0, Inches(1.8), W, Inches(0.4),
            "SECTION 2", size=16, bold=True, color=RGBColor(0x66, 0x88, 0xbb), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(2.5), W, Inches(0.8),
            "Class Design Deep Dive", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.5), W, Inches(0.4),
            "Object-oriented design for Menu, Order, Inventory, Customer, and Staff", size=18, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ======================== SLIDE 9: Menu Class ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Class: Menu & MenuItem", "backend/menu.py — Menu management with daily specials")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.3),
            "Design Decisions", size=18, bold=True, color=BLUE_DARK)
add_bullet_list(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3.5), [
    "MenuItem is a simple data class — name, price, category",
    "Menu loads all items from Menu.xlsx on initialization",
    "Daily Special randomly chosen via random.choice() each session, with configurable discount rate (default 20%)",
    "Fuzzy search: get_item() ignores case and spaces, so \"fruitsalad\" matches \"Fruit Salad\"",
], size=14, color=GRAY_700)

add_textbox(slide, Inches(6.6), Inches(1.4), Inches(5), Inches(0.3),
            "KEY METHODS", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(1.7), Inches(6), Inches(4.8),
"""class MenuItem:
    def __init__(self, name, price, category):
        self.name = name
        self.price = price
        self.category = category

class Menu:
    def __init__(self, filename="Menu.xlsx",
                 rate=0.2):
        self.menu = self.load_from_excel()
        self.special_item = self.daily_special()
        self.special_discounted_rate = rate

    def daily_special(self):
        return random.choice(self.menu)

    def get_item(self, name):
        for item in self.menu:
            if item.name.lower().replace(" ","")
               == name.lower().replace(" ",""):
                return item
        return None

    def load_from_excel(self):
        df = pd.read_excel('data/' + self.filename)
        menu = []
        for _, row in df.iterrows():
            menu.append(MenuItem(
                row["item_name"],
                row["price"],
                row["category"]))
        return menu""", font_size=10)


# ======================== SLIDE 10: Order Class ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Class: Order", "backend/order.py — Cart management, pricing, and tax calculation")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.3),
            "Design Decisions", size=18, bold=True, color=BLUE_DARK)
add_bullet_list(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3.5), [
    "Items stored as a list (not dict) — each instance individually, making add/remove simple",
    "Discount logic in get_total_after_discount() — only the daily special gets the discount",
    "Tax rate is a class constant (0.1 = 10%), easy to adjust",
    "save_to_excel() auto-increments order_id and appends to Orders.xlsx with timestamps",
], size=14, color=GRAY_700)

add_textbox(slide, Inches(6.6), Inches(1.4), Inches(5), Inches(0.3),
            "KEY METHODS", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(1.7), Inches(6), Inches(4.8),
"""class Order:
    def __init__(self, menu):
        self.items = []
        self.menu = menu
        self.TAX_RATE = 0.1

    def add_to_order(self, item):
        self.items.append(item)

    def remove_from_order(self, item):
        self.items.remove(item)

    def get_total_after_discount(self):
        total = 0.0
        for item in self.items:
            if item == self.menu.special_item:
                total += item.price * \\
                  (1 - self.menu.special_discounted_rate)
            else:
                total += item.price
        return total

    def get_items_num(self):
        # Returns {MenuItem: count, ...}
        item_counter = {}
        for item in self.items:
            if item not in item_counter:
                item_counter[item] = 1
            else:
                item_counter[item] += 1
        return item_counter""", font_size=10)


# ======================== SLIDE 11: Inventory Class ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Class: Inventory", "backend/inventory_ShuranRyu.py — Contributed by teammate Shuran Ryu")

# Collaboration highlight
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.4), Inches(12), Inches(0.9))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(1.5), Inches(11.4), Inches(0.25),
            "Collaborative Development", size=14, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(1.8), Inches(11.4), Inches(0.4),
            "The Inventory class was written by our classmate Shuran Ryu and merged into our codebase. This demonstrates our ability to integrate external code into a cohesive system, aligning interfaces so that main.py and other classes interact with it seamlessly.",
            size=12, color=GRAY_700)

add_textbox(slide, Inches(0.6), Inches(2.6), Inches(5.5), Inches(0.3),
            "Key Methods", size=18, bold=True, color=BLUE_DARK)
add_bullet_list(slide, Inches(0.6), Inches(3.0), Inches(5.5), Inches(3), [
    "check_stock(item) → True/False — safe guard before deducting",
    "deduct_stock(item) — decrements by 1 unit per call",
    "restore_stock(item) — increments by 1 when order item is deleted",
    "restock(name, amount) — staff can bulk-add any quantity",
    "load_from_excel() / save_to_excel() — Inventory.xlsx persistence",
], size=14, color=GRAY_700)

add_textbox(slide, Inches(6.6), Inches(2.5), Inches(5), Inches(0.3),
            "KEY METHODS", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(2.8), Inches(6), Inches(3.7),
"""class Inventory:
    def __init__(self, filename="Inventory.xlsx"):
        self.stock = self.load_from_excel()
        # stock = {"Americano": 5, ...}

    def check_stock(self, item):
        return self.stock.get(item.name, 0) > 0

    def deduct_stock(self, item):
        if self.check_stock(item):
            self.stock[item.name] -= 1

    def restore_stock(self, item):
        self.stock[item.name] += 1

    def restock(self, item_name, amount):
        self.stock[item_name] = \\
            self.stock.get(item_name, 0) + amount""", font_size=11)


# ======================== SLIDE 12: Customer Class ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Class: Customer & Customers", "backend/customer.py — Loyalty program with automatic point tracking")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.3),
            "Design Decisions", size=18, bold=True, color=BLUE_DARK)
add_bullet_list(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(3.5), [
    "Auto-generated IDs: class-level counter produces unique 9-digit IDs (e.g., 000000001)",
    "Phone as identifier: customers looked up by phone number; duplicate prevention built in",
    "Points policy: earn 1 point per $5 spent — clear, simple, configurable",
    "Customers class manages the collection — CRUD operations + Excel persistence",
    "Counter syncs with max existing ID on load to prevent duplicates",
], size=14, color=GRAY_700)

add_textbox(slide, Inches(6.6), Inches(1.4), Inches(5), Inches(0.3),
            "KEY METHODS", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(1.7), Inches(6), Inches(4.8),
"""class Customer:
    counter = 0  # Class-level auto-increment

    def __init__(self, telephone, name="",
                 points=0, id=None):
        if id is not None:
            self.customer_id = id
        else:
            Customer.counter += 1
            self.customer_id = \\
                str(Customer.counter).zfill(9)

        self.telephone = telephone
        self.points = points
        self.earn_point_policy = 5

    def earn_points(self, consuming):
        # 1 point per $5 spent
        self.points += consuming // 5

class Customers:
    def __init__(self, filename='Customers.xlsx'):
        self.customers = self.load_from_excel()

    def check_customer(self, telephone):
        for c in self.customers:
            if c.telephone == telephone:
                return True
        return False""", font_size=10)


# ======================== SLIDE 13: Staff Class ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Class: Staff & Staffs", "backend/staff.py — Secure staff authentication and management")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.3),
            "Design Decisions", size=18, bold=True, color=BLUE_DARK)
add_bullet_list(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(2.5), [
    "Same ID pattern as Customer — 9-digit auto-increment with class-level counter",
    "Password-protected: staff must enter both ID and password to access",
    "In the CLI, getpass hides password input for security",
    "Staffs collection mirrors the Customers pattern: check, get, add, delete, update + Excel I/O",
], size=14, color=GRAY_700)

# Consistent patterns box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.5), Inches(5.5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(4.6), Inches(5), Inches(0.25),
            "Consistent Patterns", size=13, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(4.9), Inches(5), Inches(0.6),
            "Customer and Staff follow the same architectural template — entity class + collection class + Excel persistence. This makes the codebase predictable and easy to extend.",
            size=12, color=GRAY_700)

add_textbox(slide, Inches(6.6), Inches(1.4), Inches(5), Inches(0.3),
            "STAFF LOGIN FLOW (MAIN.PY)", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(1.7), Inches(6), Inches(2.5),
"""def login_staff_system():
    while True:
        id = input("Enter staff id: ")

        if not staffs.check_staff(id):
            print("Invalid input")
            continue

        staff = staffs.get_staff(id)
        pwd = getpass("Enter password:")

        if pwd == staff.password:
            staff_menu()
            return""", font_size=11)

add_textbox(slide, Inches(6.6), Inches(4.4), Inches(5), Inches(0.3),
            "SESSION QUIT — SUMMARY ON EXIT", size=10, bold=True, color=BLUE)
add_code_box(slide, Inches(6.6), Inches(4.7), Inches(6), Inches(1.8),
"""def quit():
    for _, info in session_orders.items():
        session_orders_total += info["Total"]
        print(info["Items"])

    print(f"{len(session_orders)} orders, "
          f"${session_orders_total:.2f} earned.")""", font_size=11)


# ======================== SLIDE 14: Section - Data Persistence ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
add_textbox(slide, 0, Inches(1.8), W, Inches(0.4),
            "SECTION 3", size=16, bold=True, color=RGBColor(0x66, 0x88, 0xbb), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(2.5), W, Inches(0.8),
            "Data Persistence", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.5), W, Inches(0.4),
            "Excel-based tracking for orders, customers, and inventory", size=18, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ======================== SLIDE 15: Excel Data ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Excel Data Layer", "All data persisted to .xlsx files using pandas + openpyxl")

# Table
from pptx.util import Inches, Pt
table_data = [
    ["File", "Purpose", "Columns"],
    ["Menu.xlsx", "Menu items catalog", "item_name, price, category"],
    ["Orders.xlsx", "Complete order history", "order_id, total, date, time, [item quantities]"],
    ["Inventory.xlsx", "Current stock levels", "Item Name, Stock Numbers"],
    ["Customers.xlsx", "Customer profiles & loyalty points", "customer_id, name, telephone, points"],
    ["Staffs.xlsx", "Staff credentials", "staff_id, name, telephone, password"],
]
rows, cols = len(table_data), len(table_data[0])
table_shape = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12), Inches(3.0))
table = table_shape.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(4.5)
table.columns[2].width = Inches(5.0)

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(13)
                run.font.name = "Calibri"
                if r == 0:
                    run.font.bold = True
                    run.font.color.rgb = WHITE
                else:
                    run.font.color.rgb = GRAY_700
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE

# Special feature box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.8), Inches(12), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(4.9), Inches(11.4), Inches(0.3),
            "★  Special Feature: Loyalty Points Tracking", size=15, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.3), Inches(11.4), Inches(1.0),
            "Every checkout automatically identifies returning customers by phone number, accumulates their loyalty points in Customers.xlsx, and displays the updated balance on the receipt. New customers are auto-registered on first purchase with a unique 9-digit ID. All data persists across sessions — no data is lost when the program restarts.",
            size=14, color=GRAY_700)


# ======================== SLIDE 16: Section - Frontend ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
add_textbox(slide, 0, Inches(1.8), W, Inches(0.4),
            "SECTION 4", size=16, bold=True, color=RGBColor(0x66, 0x88, 0xbb), align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(2.5), W, Inches(0.8),
            "Frontend Application", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.5), W, Inches(0.4),
            "A React web interface built on top of the same backend", size=18, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# ======================== SLIDE 17: Frontend Overview ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "React Frontend", "Modern web interface for the Campus Cafe POS")

add_textbox(slide, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
            "To make our project more accessible and visually engaging, we built a full web application using React. A FastAPI server (app.py) wraps the exact same backend classes, exposing them as REST API endpoints.",
            size=14, color=GRAY_700)

cw2 = Inches(3.7)
ch2 = Inches(2.1)
add_info_card(slide, Inches(0.6), Inches(2.2), cw2, ch2, "Tech Stack",
              ["React 18 + Vite", "React Router v6", "FastAPI (Python)", "Vite proxy for API calls"])
add_info_card(slide, Inches(4.6), Inches(2.2), cw2, ch2, "Customer Pages",
              ["Home — landing page", "Menu — browse by category", "Cart — review & edit quantities", "Checkout — pay & get receipt"])
add_info_card(slide, Inches(8.6), Inches(2.2), cw2, ch2, "Staff Pages",
              ["Staff Login — secure auth", "Inventory Management", "Restock items in bulk", "Session summary on quit"])

# Same backend box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.6), Inches(12), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(4.7), Inches(11.4), Inches(0.3),
            "Same Backend, Two Interfaces", size=14, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.05), Inches(11.4), Inches(0.5),
            "Both the terminal CLI and the React frontend call the same Python classes. The web app adds a visual layer on top without duplicating any business logic. This demonstrates clean separation of concerns — the backend handles data and rules, while the frontend handles presentation.",
            size=13, color=GRAY_700)


# ======================== SLIDE 18: Frontend Features ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Frontend — Key Features", "Interactive UI with real-time stock validation")

cw3 = Inches(5.8)
ch3 = Inches(2.2)
add_info_card(slide, Inches(0.6), Inches(1.5), cw3, ch3, "Menu Page",
              ["4 category tabs (Salads, Drinks, Desserts, Sandwiches)",
               "Daily special banner with discount badge",
               "Inline quantity controls (+/-) with stock validation",
               "Toast notifications when stock is insufficient"])
add_info_card(slide, Inches(6.8), Inches(1.5), cw3, ch3, "Cart Page",
              ["Editable quantities (type or +/- buttons)",
               "Real-time price summary with tax & discount",
               "Stock validation on quantity increase",
               "Clear cart functionality"])
add_info_card(slide, Inches(0.6), Inches(4.0), cw3, ch3, "Checkout & Receipt",
              ["Phone number entry for loyalty points",
               "Full receipt with itemized breakdown",
               "Auto-redirect to home after 8 seconds",
               "Automatic customer registration"])
add_info_card(slide, Inches(6.8), Inches(4.0), cw3, ch3, "Staff Dashboard",
              ["Inventory table with low-stock highlighting (red)",
               "Restock form with dropdown item selection",
               "Navbar \"Quit\" button triggers session summary",
               "Modal shows all orders with item details + revenue"])


# ======================== SLIDE 19: Summary ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide_header(slide, "Summary & Key Takeaways")

cw4 = Inches(5.8)
ch4 = Inches(1.7)
add_info_card(slide, Inches(0.6), Inches(1.5), cw4, ch4, "Object-Oriented Design",
              ["5 backend classes with clear responsibilities",
               "Entity + Collection pattern for Customer and Staff",
               "Consistent Excel I/O interface across all classes"])
add_info_card(slide, Inches(6.8), Inches(1.5), cw4, ch4, "Collaborative Development",
              ["Integrated classmate's Inventory class seamlessly",
               "Demonstrates real-world teamwork and code merging",
               "Aligned interfaces for smooth integration"])
add_info_card(slide, Inches(0.6), Inches(3.5), cw4, ch4, "Full-Stack Capability",
              ["Terminal CLI + React web app sharing same backend",
               "FastAPI bridges Python classes to a modern web UI",
               "Clean separation of concerns"])
add_info_card(slide, Inches(6.8), Inches(3.5), cw4, ch4, "Data Persistence",
              ["Excel-based storage for all entities",
               "Loyalty points accumulate across sessions",
               "Complete order history with timestamps"])

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.5), Inches(12), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xdb, 0xea, 0xfe)
shape.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(5.6), Inches(11.4), Inches(0.25),
            "What We Learned", size=14, bold=True, color=BLUE_DARK)
add_textbox(slide, Inches(0.9), Inches(5.9), Inches(11.4), Inches(0.6),
            "This project reinforced OOP principles (encapsulation, separation of concerns), data persistence with pandas, input validation patterns, and how to extend a CLI application into a full-stack web application without rewriting business logic.",
            size=13, color=GRAY_700)


# ======================== SLIDE 20: Thank You ========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg_rect(slide, 0, 0, W, H, BLUE_DARK)
add_bg_rect(slide, Inches(6), 0, Inches(7.333), H, BLUE)
add_textbox(slide, 0, Inches(2.2), W, Inches(1.0),
            "Thank You", size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, 0, Inches(3.3), W, Inches(0.5),
            "Questions & Live Demo", size=22, color=RGBColor(0xcc, 0xdd, 0xee), align=PP_ALIGN.CENTER)
line_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.03))
line_shape.fill.solid()
line_shape.fill.fore_color.rgb = RGBColor(0x66, 0x88, 0xbb)
line_shape.line.fill.background()
add_textbox(slide, 0, Inches(4.5), W, Inches(0.8),
            "Campus Cafe — NEU 5001 Final Project\nPython  •  FastAPI  •  React  •  Excel",
            size=15, color=RGBColor(0x99, 0xbb, 0xdd), align=PP_ALIGN.CENTER)


# Save
output_path = "presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
